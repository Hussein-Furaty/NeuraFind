from pathlib import Path

from pptx import Presentation


class PptxParser:
    """Extracts text content from PPTX presentations."""

    def extract_text(self, file_path: str | Path) -> str:
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"Presentation does not exist: {path}")

        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Unsupported file type: {path.suffix}")

        presentation = Presentation(path)

        text_parts: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_parts.append(f"Slide {slide_index}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

        return "\n".join(text_parts)